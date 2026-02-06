#!/usr/bin/env python3
"""
Generate 5 sophisticated AVEMO presentations with proper branding and liquid glass effects
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap, qn
from pptx.oxml import parse_xml
import os

# AVEMO Brand Colors
AVEMO_ORANGE = RGBColor(0xFF, 0x79, 0x32)  # #ff7932
AVEMO_BLACK = RGBColor(0x1A, 0x1A, 0x1A)   # #1a1a1a
AVEMO_WHITE = RGBColor(0xFF, 0xFF, 0xFF)   # #ffffff
AVEMO_GRAY_DARK = RGBColor(0x33, 0x33, 0x33)  # #333333
AVEMO_GRAY_MID = RGBColor(0x66, 0x66, 0x66)   # #666666
AVEMO_GRAY_LIGHT = RGBColor(0xF5, 0xF5, 0xF5) # #f5f5f5

# Slide dimensions (16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def add_gradient_background(slide, color1, color2, direction="vertical"):
    """Add a gradient background to a slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color1
    # Note: python-pptx has limited gradient support, we'll use solid fills for now


def add_shape_with_opacity(slide, shape_type, left, top, width, height, fill_color, opacity=1.0, line_color=None):
    """Add a shape with optional transparency (simulated via overlay)"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size, font_color, bold=False, font_name="Arial", align=PP_ALIGN.LEFT):
    """Add a text box with specified styling"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def create_presentation_1():
    """
    DESIGN 1: Cinematic Dark with Orange Glow
    - Deep black background with subtle gradient
    - Large cinematic typography
    - Orange accent glows and highlights
    - Glassmorphism cards with blur effect simulation
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Dark background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x0D, 0x0D, 0x0D)
    bg.line.fill.background()
    
    # Orange gradient bar at top
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(0.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = AVEMO_ORANGE
    top_bar.line.fill.background()
    
    # Subtle glow circle behind text
    glow = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8), Inches(1), Inches(6), Inches(6))
    glow.fill.solid()
    glow.fill.fore_color.rgb = RGBColor(0xFF, 0x79, 0x32)
    # Set transparency by modifying the fill
    from pptx.oxml.ns import nsmap
    spPr = glow._element.spPr
    solidFill = spPr.find(qn('a:solidFill'))
    if solidFill is not None:
        srgbClr = solidFill.find(qn('a:srgbClr'))
        if srgbClr is not None:
            alpha = parse_xml(f'<a:alpha xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="15000"/>')
            srgbClr.append(alpha)
    glow.line.fill.background()
    
    # Main title - massive
    title = add_text_box(slide, Inches(0.8), Inches(2.2), Inches(10), Inches(1.5), 
                        "Fahrersoftware", Pt(96), AVEMO_WHITE, bold=True)
    title.text_frame.paragraphs[0].font.name = "Arial Black"
    
    # Subtitle with orange accent
    subtitle = add_text_box(slide, Inches(0.8), Inches(4), Inches(8), Inches(0.8),
                           "Die Zukunft der Fahrzeugdisposition", Pt(32), AVEMO_ORANGE, bold=True)
    
    # Glass card effect
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.2), Inches(7), Inches(1.5))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
    card.line.color.rgb = RGBColor(0xFF, 0x79, 0x32)
    card.line.width = Pt(2)
    # Add transparency
    spPr = card._element.spPr
    solidFill = spPr.find(qn('a:solidFill'))
    if solidFill is not None:
        srgbClr = solidFill.find(qn('a:srgbClr'))
        if srgbClr is not None:
            alpha = parse_xml(f'<a:alpha xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="70000"/>')
            srgbClr.append(alpha)
    
    card_text = add_text_box(slide, Inches(1), Inches(5.5), Inches(6.6), Inches(1),
                            "Eine Software für alle Standorte", Pt(20), AVEMO_WHITE)
    
    # Slide 2: Problem - Cinematic split
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x0D, 0x0D, 0x0D)
    bg.line.fill.background()
    
    # Left panel - dark
    left_panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(5), SLIDE_HEIGHT)
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = RGBColor(0x14, 0x14, 0x14)
    left_panel.line.fill.background()
    
    # Section title
    add_text_box(slide, Inches(0.5), Inches(0.8), Inches(4), Inches(0.6),
                "DAS PROBLEM", Pt(14), AVEMO_ORANGE, bold=True)
    
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(4), Inches(1),
                "Aktuelle Herausforderungen", Pt(36), AVEMO_WHITE, bold=True)
    
    # Problem cards with glass effect
    problems = [
        ("Isolierte Standorte", "Keine Vernetzung zwischen den Niederlassungen"),
        ("Manuelle Prozesse", "Excel, Outlook & Kalender führen zu Chaos"),
        ("Fehlende Daten", "Keine KPIs, keine Optimierung möglich"),
        ("Hohe Leerlaufzeit", "2,5 Stunden täglich pro Fahrzeug")
    ]
    
    for i, (title, desc) in enumerate(problems):
        y_pos = Inches(3 + i * 1)
        # Glass card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), y_pos, Inches(4), Inches(0.85))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x22)
        card.line.color.rgb = RGBColor(0xFF, 0x79, 0x32)
        card.line.width = Pt(1)
        spPr = card._element.spPr
        solidFill = spPr.find(qn('a:solidFill'))
        if solidFill is not None:
            srgbClr = solidFill.find(qn('a:srgbClr'))
            if srgbClr is not None:
                alpha = parse_xml(f'<a:alpha xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="60000"/>')
                srgbClr.append(alpha)
        
        add_text_box(slide, Inches(0.7), y_pos + Inches(0.15), Inches(3.6), Inches(0.3),
                    title, Pt(16), AVEMO_WHITE, bold=True)
        add_text_box(slide, Inches(0.7), y_pos + Inches(0.45), Inches(3.6), Inches(0.3),
                    desc, Pt(11), RGBColor(0xAA, 0xAA, 0xAA))
    
    # Right side - Large visual
    add_text_box(slide, Inches(5.5), Inches(3), Inches(7), Inches(2),
                "30 Tage", Pt(120), AVEMO_ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(5.5), Inches(5), Inches(7), Inches(0.5),
                "durchschnittliche Wartezeit auf Fahrzeuge", Pt(20), AVEMO_WHITE, align=PP_ALIGN.CENTER)
    
    # Slide 3: Solution
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x0D, 0x0D, 0x0D)
    bg.line.fill.background()
    
    # Orange diagonal accent
    diagonal = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(-2), Inches(5), Inches(20), Inches(4))
    diagonal.fill.solid()
    diagonal.fill.fore_color.rgb = RGBColor(0xFF, 0x79, 0x32)
    diagonal.rotation = -15
    diagonal.line.fill.background()
    spPr = diagonal._element.spPr
    solidFill = spPr.find(qn('a:solidFill'))
    if solidFill is not None:
        srgbClr = solidFill.find(qn('a:srgbClr'))
        if srgbClr is not None:
            alpha = parse_xml(f'<a:alpha xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="20000"/>')
            srgbClr.append(alpha)
    
    add_text_box(slide, Inches(0.8), Inches(0.8), Inches(5), Inches(0.6),
                "DIE LÖSUNG", Pt(14), AVEMO_ORANGE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(10), Inches(1),
                "Eine Software für alle Standorte", Pt(44), AVEMO_WHITE, bold=True)
    
    benefits = [
        ("Synergien nutzen", "Vernetzung aller Standorte"),
        ("Effizienz steigern", "Optimierte Fahrzeugnutzung"),
        ("Automatisierung", "Weniger manuelle Arbeit"),
        ("Kosteneinsparung", "Potenzial von 260.000 € pro Jahr")
    ]
    
    for i, (title, desc) in enumerate(benefits):
        x_pos = Inches(0.8 + (i % 2) * 6)
        y_pos = Inches(3 + (i // 2) * 1.8)
        
        # Orange accent bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_pos, y_pos, Inches(0.08), Inches(1.2))
        bar.fill.solid()
        bar.fill.fore_color.rgb = AVEMO_ORANGE
        bar.line.fill.background()
        
        add_text_box(slide, x_pos + Inches(0.2), y_pos + Inches(0.1), Inches(5), Inches(0.4),
                    title, Pt(22), AVEMO_WHITE, bold=True)
        add_text_box(slide, x_pos + Inches(0.2), y_pos + Inches(0.55), Inches(5), Inches(0.5),
                    desc, Pt(14), RGBColor(0xBB, 0xBB, 0xBB))
    
    # Slide 4: MVP Features
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x0D, 0x0D, 0x0D)
    bg.line.fill.background()
    
    add_text_box(slide, Inches(0.8), Inches(0.6), Inches(6), Inches(0.5),
                "MINIMUM VIABLE PRODUCT", Pt(12), AVEMO_ORANGE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(10), Inches(0.8),
                "Die Kernfunktionen", Pt(40), AVEMO_WHITE, bold=True)
    
    features = [
        "Optimierter Planungs- algorithmus",
        "Intuitive Driver App",
        "Modernes Web- Backoffice",
        "Integration externer Dienstleister",
        "Live-Monitoring Dashboard",
        "Automatisiertes Reporting"
    ]
    
    for i, feature in enumerate(features):
        col = i % 3
        row = i // 3
        x_pos = Inches(0.8 + col * 4)
        y_pos = Inches(2.5 + row * 2.2)
        
        # Glass card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, y_pos, Inches(3.6), Inches(1.8))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
        card.line.color.rgb = RGBColor(0xFF, 0x79, 0x32)
        card.line.width = Pt(1.5)
        spPr = card._element.spPr
        solidFill = spPr.find(qn('a:solidFill'))
        if solidFill is not None:
            srgbClr = solidFill.find(qn('a:srgbClr'))
            if srgbClr is not None:
                alpha = parse_xml(f'<a:alpha xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="50000"/>')
                srgbClr.append(alpha)
        
        # Number
        num = add_text_box(slide, x_pos + Inches(0.2), y_pos + Inches(0.15), Inches(0.5), Inches(0.5),
                          f"0{i+1}", Pt(28), AVEMO_ORANGE, bold=True)
        
        # Text
        add_text_box(slide, x_pos + Inches(0.2), y_pos + Inches(0.7), Inches(3.2), Inches(0.9),
                    feature, Pt(18), AVEMO_WHITE, bold=True)
    
    # Slide 5: Closing
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x0D, 0x0D, 0x0D)
    bg.line.fill.background()
    
    # Large orange circle background
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7), Inches(-2), Inches(10), Inches(12))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0xFF, 0x79, 0x32)
    circle.line.fill.background()
    spPr = circle._element.spPr
    solidFill = spPr.find(qn('a:solidFill'))
    if solidFill is not None:
        srgbClr = solidFill.find(qn('a:srgbClr'))
        if srgbClr is not None:
            alpha = parse_xml(f'<a:alpha xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="20000"/>')
            srgbClr.append(alpha)
    
    add_text_box(slide, Inches(0.8), Inches(2.5), Inches(10), Inches(1.2),
                "Bereit für die", Pt(48), AVEMO_WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(3.5), Inches(10), Inches(1.5),
                "Zukunft?", Pt(96), AVEMO_ORANGE, bold=True)
    
    add_text_box(slide, Inches(0.8), Inches(5.5), Inches(6), Inches(1),
                "Gemeinsam gestalten wir die digitale Transformation Ihrer Fahrzeugdisposition.",
                Pt(18), RGBColor(0xCC, 0xCC, 0xCC))
    
    # Orange accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(6.8), Inches(2), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = AVEMO_ORANGE
    line.line.fill.background()
    
    prs.save('/root/fahrerplanung/presentations/design-1-cinematic-dark.pptx')
    print("Created: design-1-cinematic-dark.pptx")


def create_presentation_2():
    """
    DESIGN 2: Clean Minimal with Liquid Glass
    - Pure white/light gray background
    - Floating glass cards with blur edges
    - Orange as sharp accent only
    - Modern Swiss typography
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = AVEMO_WHITE
    bg.line.fill.background()
    
    # Subtle gray shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(-1), Inches(-1), Inches(8), Inches(10))
    shape.fill.solid()
    shape.fill.fore_color.rgb = AVEMO_GRAY_LIGHT
    shape.line.fill.background()
    
    # Glass card for title
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), Inches(8), Inches(3.5))
    card.fill.solid()
    card.fill.fore_color.rgb = AVEMO_WHITE
    card.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    card.line.width = Pt(1)
    spPr = card._element.spPr
    solidFill = spPr.find(qn('a:solidFill'))
    if solidFill is not None:
        srgbClr = solidFill.find(qn('a:srgbClr'))
        if srgbClr is not None:
            alpha = parse_xml(f'<a:alpha xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="85000"/>')
            srgbClr.append(alpha)
    
    # Orange accent dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), Inches(2.5), Inches(0.3), Inches(0.3))
    dot.fill.solid()
    dot.fill.fore_color.rgb = AVEMO_ORANGE
    dot.line.fill.background()
    
    add_text_box(slide, Inches(2), Inches(2.4), Inches(6), Inches(0.5),
                "Fahrersoftware", Pt(14), AVEMO_GRAY_MID, bold=True)
    
    add_text_box(slide, Inches(1.5), Inches(3.2), Inches(7), Inches(1.2),
                "Die Zukunft der Fahrzeugdisposition", Pt(48), AVEMO_BLACK, bold=True)
    
    add_text_box(slide, Inches(1.5), Inches(4.5), Inches(6), Inches(0.6),
                "Eine Software für alle Standorte", Pt(20), AVEMO_GRAY_MID)
    
    # Orange line accent
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(5.3), Inches(1.5), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = AVEMO_ORANGE
    line.line.fill.background()
    
    # Slide 2: Problem
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = AVEMO_WHITE
    bg.line.fill.background()
    
    # Header
    add_text_box(slide, Inches(0.8), Inches(0.6), Inches(4), Inches(0.4),
                "Aktuelle Situation", Pt(32), AVEMO_BLACK, bold=True)
    
    # Orange underline
    underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(1.2), Inches(0.06))
    underline.fill.solid()
    underline.fill.fore_color.rgb = AVEMO_ORANGE
    underline.line.fill.background()
    
    problems = [
        ("Isolierte Standorte", "Keine Vernetzung zwischen den Niederlassungen führt zu Ineffizienz"),
        ("Manuelle Prozesse", "Abhängig von Excel, Outlook und Kalendern"),
        ("Fehlende Transparenz", "Keine KPIs für datengesteuerte Entscheidungen"),
        ("Hohe Leerlaufzeiten", "Durchschnittlich 2,5 Stunden pro Fahrzeug täglich"),
        ("Lange Wartezeiten", "30 Tage durchschnittlich auf Fahrzeuge warten")
    ]
    
    for i, (title, desc) in enumerate(problems):
        y_pos = Inches(1.6 + i * 1.1)
        
        # Number with orange
        num_bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y_pos, Inches(0.5), Inches(0.5))
        num_bg.fill.solid()
        num_bg.fill.fore_color.rgb = AVEMO_ORANGE if i == 0 else AVEMO_GRAY_LIGHT
        num_bg.line.fill.background()
        
        add_text_box(slide, Inches(0.95), y_pos + Inches(0.08), Inches(0.3), Inches(0.3),
                    str(i+1), Pt(14), AVEMO_WHITE if i == 0 else AVEMO_GRAY_MID, bold=True, align=PP_ALIGN.CENTER)
        
        # Glass card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), y_pos - Inches(0.1), Inches(11), Inches(0.9))
        card.fill.solid()
        card.fill.fore_color.rgb = AVEMO_WHITE
        card.line.color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
        card.line.width = Pt(1)
        spPr = card._element.spPr
        solidFill = spPr.find(qn('a:solidFill'))
        if solidFill is not None:
            srgbClr = solidFill.find(qn('a:srgbClr'))
            if srgbClr is not None:
                alpha = parse_xml(f'<a:alpha xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="90000"/>')
                srgbClr.append(alpha)
        
        add_text_box(slide, Inches(1.7), y_pos + Inches(0.05), Inches(3), Inches(0.4),
                    title, Pt(18), AVEMO_BLACK, bold=True)
        add_text_box(slide, Inches(1.7), y_pos + Inches(0.4), Inches(10), Inches(0.4),
                    desc, Pt(13), AVEMO_GRAY_MID)
    
    # Slide 3: Solution
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = AVEMO_WHITE
    bg.line.fill.background()
    
    # Large orange shape
    orange_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(0), Inches(6), SLIDE_HEIGHT)
    orange_shape.fill.solid()
    orange_shape.fill.fore_color.rgb = AVEMO_ORANGE
    orange_shape.line.fill.background()
    
    add_text_box(slide, Inches(0.8), Inches(0.8), Inches(6), Inches(0.5),
                "Unsere Lösung", Pt(14), AVEMO_ORANGE, bold=True)
    
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(6), Inches(1),
                "Eine Software für alle Standorte", Pt(36), AVEMO_BLACK, bold=True)
    
    solutions = [
        ("Synergien", "Alle Standorte vernetzt"),
        ("Effizienz", "Optimierte Prozesse"),
        ("Automatisierung", "Weniger manuelle Arbeit"),
        ("Einsparung", "260.000 € jährlich")
    ]
    
    for i, (title, desc) in enumerate(solutions):
        y_pos = Inches(3 + i * 1.1)
        
        # Icon placeholder (circle)
        icon = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y_pos, Inches(0.6), Inches(0.6))
        icon.fill.solid()
        icon.fill.fore_color.rgb = AVEMO_ORANGE
        icon.line.fill.background()
        
        add_text_box(slide, Inches(1.6), y_pos + Inches(0.1), Inches(2.5), Inches(0.4),
                    title, Pt(22), AVEMO_BLACK, bold=True)
        add_text_box(slide, Inches(1.6), y_pos + Inches(0.5), Inches(5), Inches(0.4),
                    desc, Pt(14), AVEMO_GRAY_MID)
    
    # White text on orange
    add_text_box(slide, Inches(8.5), Inches(3), Inches(4), Inches(3),
                "Transformieren Sie Ihre Fahrzeugdisposition mit moderner Technologie",
                Pt(24), AVEMO_WHITE, bold=True)
    
    # Slide 4: MVP
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = AVEMO_WHITE
    bg.line.fill.background()
    
    add_text_box(slide, Inches(0.8), Inches(0.6), Inches(6), Inches(0.4),
                "Minimum Viable Product", Pt(14), AVEMO_ORANGE, bold=True)
    
    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(10), Inches(0.8),
                "Die Kernfunktionen", Pt(40), AVEMO_BLACK, bold=True)
    
    features = [
        "Planungsalgorithmus",
        "Driver App",
        "Web-Backoffice",
        "Externe Dienstleister",
        "Monitoring",
        "Reporting"
    ]
    
    for i, feature in enumerate(features):
        row = i // 3
        col = i % 3
        x = Inches(0.8 + col * 4.1)
        y = Inches(2.5 + row * 2.2)
        
        # Glass card with subtle shadow effect
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.8), Inches(1.9))
        card.fill.solid()
        card.fill.fore_color.rgb = AVEMO_WHITE
        card.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        card.line.width = Pt(1)
        
        # Orange accent bar at top
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(3.8), Inches(0.08))
        bar.fill.solid()
        bar.fill.fore_color.rgb = AVEMO_ORANGE
        bar.line.fill.background()
        
        add_text_box(slide, x + Inches(0.3), y + Inches(0.5), Inches(3.2), Inches(1),
                    feature, Pt(20), AVEMO_BLACK, bold=True)
        
        # Small number
        add_text_box(slide, x + Inches(0.3), y + Inches(0.25), Inches(0.5), Inches(0.3),
                    f"0{i+1}", Pt(12), AVEMO_ORANGE, bold=True)
    
    # Slide 5: Closing
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = AVEMO_WHITE
    bg.line.fill.background()
    
    # Subtle pattern of circles
    for pos in [(2, 1), (10, 5), (11, 0.5), (1, 6)]:
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(pos[0]), Inches(pos[1]), Inches(0.3), Inches(0.3))
        circle.fill.solid()
        circle.fill.fore_color.rgb = AVEMO_ORANGE
        circle.line.fill.background()
        spPr = circle._element.spPr
        solidFill = spPr.find(qn('a:solidFill'))
        if solidFill is not None:
            srgbClr = solidFill.find(qn('a:srgbClr'))
            if srgbClr is not None:
                alpha = parse_xml(f'<a:alpha xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="30000"/>')
                srgbClr.append(alpha)
    
    add_text_box(slide, Inches(0.8), Inches(2.8), Inches(10), Inches(1),
                "Bereit für die", Pt(36), AVEMO_GRAY_MID)
    add_text_box(slide, Inches(0.8), Inches(3.5), Inches(10), Inches(1.2),
                "Zukunft?", Pt(72), AVEMO_BLACK, bold=True)
    
    # Orange underline
    underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.7), Inches(3), Inches(0.08))
    underline.fill.solid()
    underline.fill.fore_color.rgb = AVEMO_ORANGE
    underline.line.fill.background()
    
    add_text_box(slide, Inches(0.8), Inches(5.2), Inches(8), Inches(0.8),
                "Gemeinsam gestalten wir die digitale Transformation.",
                Pt(18), AVEMO_GRAY_MID)
    
    prs.save('/root/fahrerplanung/presentations/design-2-minimal-glass.pptx')
    print("Created: design-2-minimal-glass.pptx")


def create_presentation_3():
    """
    DESIGN 3: Bold Typographic with Orange Blocks
    - High contrast black/white/orange
    - Massive typography as design element
    - Geometric orange blocks and shapes
    - Editorial magazine style
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = AVEMO_BLACK
    bg.line.fill.background()
    
    # Large typography
    add_text_box(slide, Inches(-0.2), Inches(1.5), Inches(14), Inches(2),
                "FAHRER", Pt(140), AVEMO_WHITE, bold=True)
    
    # Orange block behind "SOFTWARE"
    orange_block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.3), Inches(8), Inches(1.8))
    orange_block.fill.solid()
    orange_block.fill.fore_color.rgb = AVEMO_ORANGE
    orange_block.line.fill.background()
    
    add_text_box(slide, Inches(5.7), Inches(3.4), Inches(7), Inches(1.6),
                "SOFTWARE", Pt(100), AVEMO_BLACK, bold=True)
    
    # Subtitle
    add_text_box(slide, Inches(0.8), Inches(5.8), Inches(8), Inches(0.6),
                "Die Zukunft der Fahrzeugdisposition", Pt(24), AVEMO_WHITE)
    
    # Slide 2: Problem
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Split layout
    left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(4), SLIDE_HEIGHT)
    left.fill.solid()
    left.fill.fore_color.rgb = AVEMO_ORANGE
    left.line.fill.background()
    
    right = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), 0, Inches(9.333), SLIDE_HEIGHT)
    right.fill.solid()
    right.fill.fore_color.rgb = AVEMO_WHITE
    right.line.fill.background()
    
    # Left side - large number
    add_text_box(slide, Inches(0.5), Inches(2.5), Inches(3), Inches(2),
                "01", Pt(120), AVEMO_BLACK, bold=True)
    
    add_text_box(slide, Inches(0.8), Inches(5), Inches(3), Inches(1),
                "PROBLEM", Pt(18), AVEMO_BLACK, bold=True)
    
    # Right side content
    add_text_box(slide, Inches(4.5), Inches(0.8), Inches(8), Inches(0.6),
                "Aktuelle Herausforderungen", Pt(32), AVEMO_BLACK, bold=True)
    
    problems = [
        "Isolierte Standorte",
        "Manuelle Prozesse",
        "Fehlende KPIs",
        "2,5h Leerlauf",
        "30 Tage Wartezeit"
    ]
    
    for i, problem in enumerate(problems):
        y = Inches(1.8 + i * 1)
        
        # Black accent line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), y + Inches(0.15), Inches(0.4), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = AVEMO_BLACK
        line.line.fill.background()
        
        add_text_box(slide, Inches(5.1), y, Inches(7), Inches(0.5),
                    problem, Pt(22), AVEMO_BLACK, bold=True)
    
    # Slide 3: Solution
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = AVEMO_BLACK
    bg.line.fill.background()
    
    # White block
    white_block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(1.5), Inches(9), Inches(4.5))
    white_block.fill.solid()
    white_block.fill.fore_color.rgb = AVEMO_WHITE
    white_block.line.fill.background()
    
    add_text_box(slide, Inches(2.5), Inches(1.8), Inches(8), Inches(0.5),
                "DIE LÖSUNG", Pt(14), AVEMO_ORANGE, bold=True)
    
    add_text_box(slide, Inches(2.5), Inches(2.4), Inches(8), Inches(1),
                "Eine Software für alle Standorte", Pt(36), AVEMO_BLACK, bold=True)
    
    solutions = [
        ("Synergien", "Vernetzung aller Standorte"),
        ("Effizienz", "Optimierte Fahrzeugnutzung"),
        ("Automatisierung", "Reduzierung manueller Arbeit"),
        ("Kosten", "260.000 € Einsparung")
    ]
    
    for i, (title, desc) in enumerate(solutions):
        y = Inches(3.6 + i * 0.6)
        add_text_box(slide, Inches(2.5), y, Inches(3), Inches(0.4),
                    title, Pt(16), AVEMO_ORANGE, bold=True)
        add_text_box(slide, Inches(5), y, Inches(5), Inches(0.4),
                    desc, Pt(14), AVEMO_GRAY_MID)
    
    # Slide 4: MVP Features
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = AVEMO_WHITE
    bg.line.fill.background()
    
    # Large background text
    add_text_box(slide, Inches(-1), Inches(5.5), Inches(20), Inches(2),
                "MVP", Pt(200), RGBColor(0xF0, 0xF0, 0xF0), bold=True)
    
    add_text_box(slide, Inches(0.8), Inches(0.6), Inches(6), Inches(0.5),
                "MINIMUM VIABLE PRODUCT", Pt(14), AVEMO_ORANGE, bold=True)
    
    features = [
        ("Planungs- algorithmus", "Optimierte Tourenplanung"),
        ("Driver App", "Intuitive Bedienung"),
        ("Web- Backoffice", "Moderne Oberfläche"),
        ("Monitoring", "Echtzeit-Übersicht"),
        ("Reporting", "Automatisierte Reports"),
        ("Integration", "Externe Dienstleister")
    ]
    
    for i, (title, desc) in enumerate(features):
        col = i % 3
        row = i // 3
        x = Inches(0.8 + col * 4.1)
        y = Inches(1.5 + row * 2.8)
        
        # Number
        add_text_box(slide, x, y, Inches(1), Inches(0.8),
                    f"0{i+1}", Pt(48), AVEMO_ORANGE, bold=True)
        
        add_text_box(slide, x, y + Inches(0.8), Inches(3.8), Inches(0.6),
                    title, Pt(20), AVEMO_BLACK, bold=True)
        add_text_box(slide, x, y + Inches(1.3), Inches(3.8), Inches(0.5),
                    desc, Pt(13), AVEMO_GRAY_MID)
    
    # Slide 5: Closing
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = AVEMO_BLACK
    bg.line.fill.background()
    
    # Orange diagonal
    diag = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6), Inches(-2), Inches(4), Inches(12))
    diag.fill.solid()
    diag.fill.fore_color.rgb = AVEMO_ORANGE
    diag.rotation = -25
    diag.line.fill.background()
    
    add_text_box(slide, Inches(0.8), Inches(2.5), Inches(8), Inches(1),
                "Bereit für", Pt(36), AVEMO_WHITE)
    add_text_box(slide, Inches(0.8), Inches(3.3), Inches(8), Inches(1.2),
                "die Zukunft?", Pt(60), AVEMO_ORANGE, bold=True)
    
    add_text_box(slide, Inches(0.8), Inches(4.8), Inches(7), Inches(1),
                "Gemeinsam transformieren wir Ihre Fahrzeugdisposition.",
                Pt(16), AVEMO_GRAY_MID)
    
    prs.save('/root/fahrerplanung/presentations/design-3-bold-editorial.pptx')
    print("Created: design-3-bold-editorial.pptx")


def create_presentation_4():
    """
    DESIGN 4: Modern Asymmetric Editorial
    - Asymmetric layouts
    - Image placeholders with overlays
    - Soft gradients and rounded corners
    - Contemporary web aesthetic
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    # Slide 1: Title with asymmetric layout
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Soft gradient background (using solid color)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA)
    bg.line.fill.background()
    
    # Asymmetric orange shape
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(-1), Inches(7), Inches(10))
    shape.fill.solid()
    shape.fill.fore_color.rgb = AVEMO_ORANGE
    shape.line.fill.background()
    shape.adjustments[0] = 0.1
    
    # Content on white side
    add_text_box(slide, Inches(1), Inches(2.5), Inches(5), Inches(1),
                "Fahrersoftware", Pt(52), AVEMO_BLACK, bold=True)
    
    add_text_box(slide, Inches(1), Inches(3.6), Inches(5), Inches(0.8),
                "Die Zukunft der Fahrzeugdisposition", Pt(20), AVEMO_GRAY_MID)
    
    # Subtle line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(4.6), Inches(2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = AVEMO_ORANGE
    line.line.fill.background()
    
    # Content on orange side
    add_text_box(slide, Inches(7.5), Inches(3), Inches(5), Inches(2),
                "Eine Software für alle Standorte",
                Pt(28), AVEMO_WHITE, bold=True)
    
    # Slide 2: Problem
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = AVEMO_WHITE
    bg.line.fill.background()
    
    # Header with orange accent
    header_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(1.5))
    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = AVEMO_ORANGE
    header_bg.line.fill.background()
    
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
                "Aktuelle Herausforderungen", Pt(32), AVEMO_WHITE, bold=True)
    
    # Problem cards in asymmetric grid
    problems = [
        ("Isolierte Standorte", "Keine Vernetzung zwischen den Niederlassungen", 0.8, 1.8),
        ("Manuelle Prozesse", "Abhängigkeit von Excel und Outlook", 4.5, 1.8),
        ("Fehlende Daten", "Keine KPIs für Entscheidungen", 8.2, 1.8),
        ("Hohe Leerlaufzeit", "2,5 Stunden pro Fahrzeug täglich", 2.6, 4.2),
        ("Lange Wartezeiten", "30 Tage durchschnittlich", 6.3, 4.2),
    ]
    
    for title, desc, x, y in problems:
        # Card with rounded corners
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.5), Inches(2))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA)
        card.line.fill.background()
        card.adjustments[0] = 0.15
        
        add_text_box(slide, Inches(x + 0.2), Inches(y + 0.3), Inches(3.1), Inches(0.6),
                    title, Pt(18), AVEMO_BLACK, bold=True)
        add_text_box(slide, Inches(x + 0.2), Inches(y + 1), Inches(3.1), Inches(0.8),
                    desc, Pt(13), AVEMO_GRAY_MID)
    
    # Slide 3: Solution
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA)
    bg.line.fill.background()
    
    # Large rounded card
    main_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.8), Inches(11.7), Inches(6))
    main_card.fill.solid()
    main_card.fill.fore_color.rgb = AVEMO_WHITE
    main_card.line.fill.background()
    main_card.adjustments[0] = 0.05
    
    add_text_box(slide, Inches(1.3), Inches(1.3), Inches(6), Inches(0.5),
                "Die Lösung", Pt(14), AVEMO_ORANGE, bold=True)
    
    add_text_box(slide, Inches(1.3), Inches(1.9), Inches(10), Inches(0.8),
                "Eine Software für alle Standorte", Pt(36), AVEMO_BLACK, bold=True)
    
    solutions = [
        ("Synergien nutzen", "Vernetzung aller Standorte für maximale Effizienz"),
        ("Effizienz steigern", "Optimale Nutzung jedes Fahrzeugs"),
        ("Automatisierung", "Reduzierung manueller Prozesse"),
        ("Kosteneinsparung", "Potenzial von 260.000 € jährlich")
    ]
    
    for i, (title, desc) in enumerate(solutions):
        y = Inches(3 + i * 0.9)
        
        # Icon circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.3), y, Inches(0.4), Inches(0.4))
        circle.fill.solid()
        circle.fill.fore_color.rgb = AVEMO_ORANGE
        circle.line.fill.background()
        
        add_text_box(slide, Inches(2), y, Inches(3), Inches(0.4),
                    title, Pt(18), AVEMO_BLACK, bold=True)
        add_text_box(slide, Inches(5), y, Inches(6), Inches(0.4),
                    desc, Pt(14), AVEMO_GRAY_MID)
    
    # Slide 4: MVP
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = AVEMO_WHITE
    bg.line.fill.background()
    
    add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.5),
                "Minimum Viable Product", Pt(14), AVEMO_ORANGE, bold=True)
    
    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(8), Inches(0.7),
                "Die Kernfunktionen", Pt(36), AVEMO_BLACK, bold=True)
    
    features = [
        "Planungsalgorithmus",
        "Driver App",
        "Web-Backoffice",
        "Externe Dienstleister",
        "Monitoring",
        "Reporting"
    ]
    
    # Asymmetric 2-3-1 layout
    positions = [
        (0.8, 2.2), (4.5, 2.2),
        (0.8, 4), (4.5, 4), (8.2, 4),
        (4.5, 5.8)
    ]
    
    for i, feature in enumerate(features):
        x, y = positions[i]
        
        # Rounded pill shape
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.3), Inches(1.4))
        pill.fill.solid()
        pill.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA) if i % 2 == 0 else AVEMO_ORANGE
        pill.line.fill.background()
        pill.adjustments[0] = 0.5
        
        text_color = AVEMO_BLACK if i % 2 == 0 else AVEMO_WHITE
        add_text_box(slide, Inches(x + 0.3), Inches(y + 0.45), Inches(2.7), Inches(0.6),
                    feature, Pt(18), text_color, bold=True)
    
    # Slide 5: Closing
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Gradient-like effect with two colors
    bg1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(3.75))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = AVEMO_BLACK
    bg1.line.fill.background()
    
    bg2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.75), SLIDE_WIDTH, Inches(3.75))
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = AVEMO_ORANGE
    bg2.line.fill.background()
    
    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(10), Inches(1),
                "Bereit für die", Pt(36), AVEMO_WHITE)
    add_text_box(slide, Inches(0.8), Inches(2.1), Inches(10), Inches(1),
                "Zukunft?", Pt(72), AVEMO_WHITE, bold=True)
    
    add_text_box(slide, Inches(0.8), Inches(4.5), Inches(10), Inches(1),
                "Gemeinsam gestalten wir die digitale Transformation Ihrer Fahrzeugdisposition.",
                Pt(18), AVEMO_BLACK)
    
    prs.save('/root/fahrerplanung/presentations/design-4-asymmetric-modern.pptx')
    print("Created: design-4-asymmetric-modern.pptx")


def create_presentation_5():
    """
    DESIGN 5: Data-Driven Dashboard Style
    - Dark tech aesthetic
    - Card-based layout like dashboard
    - Neon orange accents
    - Modern data visualization feel
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Dark gradient background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x0A, 0x0E, 0x17)
    bg.line.fill.background()
    
    # Grid pattern simulation
    for i in range(15):
        x = Inches(i * 0.9)
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, 0, Inches(0.01), SLIDE_HEIGHT)
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0x1A, 0x1F, 0x2E)
        line.line.fill.background()
        spPr = line._element.spPr
        solidFill = spPr.find(qn('a:solidFill'))
        if solidFill is not None:
            srgbClr = solidFill.find(qn('a:srgbClr'))
            if srgbClr is not None:
                alpha = parse_xml(f'<a:alpha xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="50000"/>')
                srgbClr.append(alpha)
    
    # Main content card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1.5), Inches(11), Inches(4.5))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x11, 0x18, 0x27)
    card.line.color.rgb = RGBColor(0xFF, 0x79, 0x32)
    card.line.width = Pt(2)
    spPr = card._element.spPr
    solidFill = spPr.find(qn('a:solidFill'))
    if solidFill is not None:
        srgbClr = solidFill.find(qn('a:srgbClr'))
        if srgbClr is not None:
            alpha = parse_xml(f'<a:alpha xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="80000"/>')
            srgbClr.append(alpha)
    
    # Orange accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.5), Inches(0.2), Inches(4.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = AVEMO_ORANGE
    bar.line.fill.background()
    
    add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
                "Fahrersoftware", Pt(64), AVEMO_WHITE, bold=True)
    
    add_text_box(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(0.6),
                "Die Zukunft der Fahrzeugdisposition", Pt(24), RGBColor(0x8A, 0x9A, 0xB0))
    
    add_text_box(slide, Inches(1.5), Inches(4.2), Inches(10), Inches(0.6),
                "Eine Software für alle Standorte", Pt(20), AVEMO_ORANGE, bold=True)
    
    # Bottom stats
    stats = [("260K€", "Einsparung"), ("30", "Tage gespart"), ("2.5h", "Weniger Leerlauf")]
    for i, (num, label) in enumerate(stats):
        x = Inches(1.5 + i * 3.5)
        add_text_box(slide, x, Inches(6.2), Inches(3), Inches(0.6),
                    num, Pt(32), AVEMO_ORANGE, bold=True)
        add_text_box(slide, x, Inches(6.8), Inches(3), Inches(0.4),
                    label, Pt(12), RGBColor(0x8A, 0x9A, 0xB0))
    
    # Slide 2: Problem
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x0A, 0x0E, 0x17)
    bg.line.fill.background()
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(0x11, 0x18, 0x27)
    header.line.fill.background()
    
    add_text_box(slide, Inches(0.8), Inches(0.35), Inches(10), Inches(0.6),
                "SYSTEM STATUS: PROBLEME ERKANNT", Pt(18), AVEMO_ORANGE, bold=True)
    
    # Problem cards in dashboard style
    problems = [
        ("Standorte", "5", "Isoliert", "CRITICAL"),
        ("Prozesse", "100%", "Manuell", "WARNING"),
        ("Leerlauf", "2.5h", "Pro Tag", "WARNING"),
        ("Wartezeit", "30", "Tage", "CRITICAL")
    ]
    
    for i, (title, value, unit, status) in enumerate(problems):
        row = i // 2
        col = i % 2
        x = Inches(0.8 + col * 6)
        y = Inches(1.6 + row * 2.8)
        
        # Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.5), Inches(2.4))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0x11, 0x18, 0x27)
        card.line.color.rgb = RGBColor(0xFF, 0x79, 0x32) if status == "CRITICAL" else RGBColor(0x33, 0x44, 0x66)
        card.line.width = Pt(1.5)
        
        add_text_box(slide, x + Inches(0.3), y + Inches(0.2), Inches(3), Inches(0.4),
                    title, Pt(14), RGBColor(0x8A, 0x9A, 0xB0))
        
        add_text_box(slide, x + Inches(0.3), y + Inches(0.7), Inches(3), Inches(0.8),
                    value, Pt(48), AVEMO_WHITE, bold=True)
        
        add_text_box(slide, x + Inches(0.3), y + Inches(1.5), Inches(3), Inches(0.4),
                    unit, Pt(14), RGBColor(0x8A, 0x9A, 0xB0))
        
        # Status indicator
        status_color = AVEMO_ORANGE if status == "CRITICAL" else RGBColor(0xFF, 0xC1, 0x07)
        indicator = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(4.5), y + Inches(0.3), Inches(0.5), Inches(0.5))
        indicator.fill.solid()
        indicator.fill.fore_color.rgb = status_color
        indicator.line.fill.background()
        
        add_text_box(slide, x + Inches(4.2), y + Inches(0.9), Inches(1), Inches(0.4),
                    status, Pt(10), status_color, bold=True)
    
    # Slide 3: Solution
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x0A, 0x0E, 0x17)
    bg.line.fill.background()
    
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(0x11, 0x18, 0x27)
    header.line.fill.background()
    
    add_text_box(slide, Inches(0.8), Inches(0.35), Inches(10), Inches(0.6),
                "SYSTEM UPGRADE: LÖSUNG BEREIT", Pt(18), AVEMO_ORANGE, bold=True)
    
    # Large feature card
    main_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(7), Inches(5.4))
    main_card.fill.solid()
    main_card.fill.fore_color.rgb = RGBColor(0x11, 0x18, 0x27)
    main_card.line.color.rgb = AVEMO_ORANGE
    main_card.line.width = Pt(2)
    
    add_text_box(slide, Inches(1.1), Inches(1.9), Inches(6), Inches(0.8),
                "Eine Software für alle Standorte", Pt(28), AVEMO_WHITE, bold=True)
    
    benefits = [
        ("Synergien", "Vernetzung aller Standorte"),
        ("Effizienz", "Optimale Nutzung"),
        ("Automatisierung", "Weniger manuelle Arbeit"),
        ("Einsparung", "260.000 € jährlich")
    ]
    
    for i, (title, desc) in enumerate(benefits):
        y = Inches(3 + i * 0.9)
        
        # Checkmark
        check = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.1), y + Inches(0.05), Inches(0.3), Inches(0.3))
        check.fill.solid()
        check.fill.fore_color.rgb = AVEMO_ORANGE
        check.line.fill.background()
        
        add_text_box(slide, Inches(1.6), y, Inches(2), Inches(0.4),
                    title, Pt(16), AVEMO_ORANGE, bold=True)
        add_text_box(slide, Inches(3.2), y, Inches(4), Inches(0.4),
                    desc, Pt(14), RGBColor(0x8A, 0x9A, 0xB0))
    
    # Side panel with metrics
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8), Inches(1.6), Inches(4.5), Inches(5.4))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(0x11, 0x18, 0x27)
    panel.line.fill.background()
    
    add_text_box(slide, Inches(8.3), Inches(1.9), Inches(4), Inches(0.5),
                "PROJEKTIERTE METRIKEN", Pt(12), AVEMO_ORANGE, bold=True)
    
    metrics = [
        ("ROI", "< 12 Monate"),
        ("Effizienz", "+ 40%"),
        ("Transparenz", "100%"),
        ("Standorte", "Verbindung")
    ]
    
    for i, (label, value) in enumerate(metrics):
        y = Inches(2.6 + i * 1.1)
        add_text_box(slide, Inches(8.3), y, Inches(2), Inches(0.4),
                    label, Pt(14), RGBColor(0x8A, 0x9A, 0xB0))
        add_text_box(slide, Inches(8.3), y + Inches(0.4), Inches(4), Inches(0.5),
                    value, Pt(24), AVEMO_WHITE, bold=True)
    
    # Slide 4: MVP
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x0A, 0x0E, 0x17)
    bg.line.fill.background()
    
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(0x11, 0x18, 0x27)
    header.line.fill.background()
    
    add_text_box(slide, Inches(0.8), Inches(0.35), Inches(10), Inches(0.6),
                "MODULE: MINIMUM VIABLE PRODUCT", Pt(18), AVEMO_ORANGE, bold=True)
    
    features = [
        ("ALGORITHM", "Planungs-Engine"),
        ("MOBILE", "Driver App"),
        ("WEB", "Backoffice"),
        ("INTEGRATION", "Externe APIs"),
        ("MONITOR", "Live-Dashboard"),
        ("ANALYTICS", "Reporting")
    ]
    
    for i, (code, name) in enumerate(features):
        col = i % 3
        row = i // 3
        x = Inches(0.8 + col * 4.1)
        y = Inches(1.6 + row * 2.8)
        
        # Module card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.8), Inches(2.4))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0x11, 0x18, 0x27)
        card.line.color.rgb = RGBColor(0x33, 0x44, 0x66)
        card.line.width = Pt(1)
        
        # Code label
        add_text_box(slide, x + Inches(0.2), y + Inches(0.2), Inches(3.4), Inches(0.35),
                    code, Pt(11), AVEMO_ORANGE, bold=True)
        
        # Name
        add_text_box(slide, x + Inches(0.2), y + Inches(0.8), Inches(3.4), Inches(0.8),
                    name, Pt(22), AVEMO_WHITE, bold=True)
        
        # Module number
        add_text_box(slide, x + Inches(3), y + Inches(1.9), Inches(0.6), Inches(0.4),
                    f"M0{i+1}", Pt(10), RGBColor(0x8A, 0x9A, 0xB0))
    
    # Slide 5: Closing
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x0A, 0x0E, 0x17)
    bg.line.fill.background()
    
    # Central card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(1.5), Inches(9.3), Inches(4.5))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x11, 0x18, 0x27)
    card.line.color.rgb = AVEMO_ORANGE
    card.line.width = Pt(3)
    
    # Animated border effect (multiple thin lines)
    for offset in [0.02, 0.04]:
        border = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                       Inches(2 - offset), Inches(1.5 - offset), 
                                       Inches(9.3 + offset*2), Inches(4.5 + offset*2))
        border.fill.background()
        border.line.color.rgb = AVEMO_ORANGE
        border.line.width = Pt(1)
        spPr = border._element.spPr
        
    add_text_box(slide, Inches(2.5), Inches(2.2), Inches(8), Inches(1),
                "SYSTEM READY", Pt(18), AVEMO_ORANGE, bold=True)
    
    add_text_box(slide, Inches(2.5), Inches(3), Inches(8), Inches(1),
                "Bereit für die Zukunft?", Pt(44), AVEMO_WHITE, bold=True)
    
    add_text_box(slide, Inches(2.5), Inches(4.2), Inches(8), Inches(0.8),
                "Gemeinsam transformieren wir Ihre Fahrzeugdisposition.",
                Pt(16), RGBColor(0x8A, 0x9A, 0xB0))
    
    prs.save('/root/fahrerplanung/presentations/design-5-dashboard-tech.pptx')
    print("Created: design-5-dashboard-tech.pptx")


if __name__ == "__main__":
    print("Creating 5 sophisticated AVEMO presentations...")
    print("="*60)
    create_presentation_1()
    create_presentation_2()
    create_presentation_3()
    create_presentation_4()
    create_presentation_5()
    print("="*60)
    print("All presentations created successfully!")
